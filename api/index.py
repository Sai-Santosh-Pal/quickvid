from flask import Flask, request, send_file, render_template_string, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests, io, re

app = Flask(__name__)

IMGBB_API_KEY = "9d1c67a0ea72097a70e086c48c838c35"  # Replace with your key

# ---------------- Helpers ----------------

def parse_attrs(block: str) -> dict:
    return {k: v.strip('"') for k, v in re.findall(r'(\w+)=(".*?"|\S+)', block)}

def fetch_image(src):
    """Return a file-like object for PPTX; accepts URL or local path"""
    if src.startswith("http://") or src.startswith("https://"):
        r = requests.get(src, stream=True)
        r.raise_for_status()
        if "image" not in r.headers.get("Content-Type",""):
            raise ValueError(f"URL is not an image: {src}")
        return io.BytesIO(r.content)
    else:
        return src  # local path

# ---------------- PPT Generator ----------------

def generate_ppt(script: str) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    M = Inches(0.5)

    def centered_text(word):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(int(W*0.25), int(H*0.35), int(W*0.5), int(H*0.3))
        p = box.text_frame.paragraphs[0]
        p.text = word
        p.font.name = "Arial"
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(0,0,0)
        p.alignment = PP_ALIGN.CENTER

    def image_only(cfg):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = fetch_image(cfg["src"])
        width_pct = float(cfg.get("width","60%").strip("%")) / 100
        width = int(W * width_pct)
        left = int((W - width) / 2)
        slide.shapes.add_picture(img, left, int(M), width=width)

    def mix(cfg):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = fetch_image(cfg["src"])
        align = cfg.get("align","left")
        img_left = int(M if align=="left" else W/2 + M)
        txt_left = int(W/2 + M if align=="left" else M)
        slide.shapes.add_picture(img, img_left, int(M), width=int(W/2 - M*2))
        box = slide.shapes.add_textbox(txt_left, int(H*0.35), int(W/2 - M*2), int(H*0.3))
        p = box.text_frame.paragraphs[0]
        p.text = cfg.get("text","")
        p.font.name = "Arial"
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0,0,0)
        p.alignment = PP_ALIGN.CENTER

    def video(cfg):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(int(W*0.25), int(H*0.35), int(W*0.5), int(H*0.3))
        p = box.text_frame.paragraphs[0]
        p.text = f"VIDEO PLACEHOLDER\n{cfg.get('src','')}"
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

    # Split script by --- for separate files
    scripts = [s.strip() for s in script.split("---") if s.strip()]
    if len(scripts) > 1:
        files = []
        for idx, s in enumerate(scripts):
            buf = io.BytesIO()
            prs_tmp = generate_single_ppt(s)
            prs_tmp.save(buf)
            buf.seek(0)
            files.append(("output_%d.pptx"%idx, buf.read()))
        return files
    else:
        return generate_single_ppt(script)

def generate_single_ppt(script: str) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    M = Inches(0.5)

    def centered_text(word):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(int(W*0.25), int(H*0.35), int(W*0.5), int(H*0.3))
        p = box.text_frame.paragraphs[0]
        p.text = word
        p.font.name = "Arial"
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(0,0,0)
        p.alignment = PP_ALIGN.CENTER

    def image_only(cfg):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = fetch_image(cfg["src"])
        width_pct = float(cfg.get("width","60%").strip("%")) / 100
        width = int(W * width_pct)
        left = int((W - width) / 2)
        slide.shapes.add_picture(img, left, int(M), width=width)

    def mix(cfg):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = fetch_image(cfg["src"])
        align = cfg.get("align","left")
        img_left = int(M if align=="left" else W/2 + M)
        txt_left = int(W/2 + M if align=="left" else M)
        slide.shapes.add_picture(img, img_left, int(M), width=int(W/2 - M*2))
        box = slide.shapes.add_textbox(txt_left, int(H*0.35), int(W/2 - M*2), int(H*0.3))
        p = box.text_frame.paragraphs[0]
        p.text = cfg.get("text","")
        p.font.name = "Arial"
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0,0,0)
        p.alignment = PP_ALIGN.CENTER

    def video(cfg):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(int(W*0.25), int(H*0.35), int(W*0.5), int(H*0.3))
        p = box.text_frame.paragraphs[0]
        p.text = f"VIDEO PLACEHOLDER\n{cfg.get('src','')}"
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

    tokens = re.findall(r'\[(IMAGE|MIX|VIDEO)(.*?)\]|(\S+)', script)
    for kind, block, word in tokens:
        if word:
            centered_text(word)
        elif kind == "IMAGE":
            image_only(parse_attrs(block))
        elif kind == "MIX":
            mix(parse_attrs(block))
        elif kind == "VIDEO":
            video(parse_attrs(block))
    return prs

# ---------------- IMGBB ----------------

def upload_to_imgbb(file):
    r = requests.post(
        "https://api.imgbb.com/1/upload",
        params={"key": IMGBB_API_KEY},
        files={"image": file.read()}
    )
    data = r.json()
    if not data.get("success"):
        raise RuntimeError("ImgBB upload failed")
    return data["data"]["url"]

# ---------------- ROUTES ----------------

@app.route("/")
def index():
    return render_template_string("""
<!doctype html>
<html>
<head>
<title>PPT Script Builder</title>
<link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.2/dist/css/bootstrap.min.css" rel="stylesheet">
<style>textarea{font-family:monospace}</style>
</head>
<body class="container mt-4">

<h3>PPT Script Builder</h3>

<textarea id="script" class="form-control mb-3" rows="10"
placeholder="Words = slides
[IMAGE src=URL width=60%]
[MIX text=&quot;Hello&quot; src=URL align=left]
[VIDEO src=URL]
Use --- for separate files"></textarea>

<button class="btn btn-primary mb-3" onclick="pick()">Add Image</button>

<form method="POST" action="/generate">
<input type="hidden" name="script" id="final">
<button class="btn btn-success">Generate PPT</button>
</form>

<input type="file" id="file" hidden>

<script>
function pick(){document.getElementById("file").click();}
document.getElementById("file").onchange=()=>{
 let f=document.getElementById("file").files[0];
 let fd=new FormData(); fd.append("file",f);
 fetch("/upload",{method:"POST",body:fd})
 .then(r=>r.json()).then(d=>{
   if(d.error){alert(d.error);return;}
   let t=document.getElementById("script");
   let tok="[IMAGE src="+d.url+" width=60%]";
   let s=t.selectionStart;
   t.value=t.value.slice(0,s)+tok+t.value.slice(s);
   t.selectionStart=t.selectionEnd=s+tok.length;
 });
}
document.querySelector("form").onsubmit=()=>{
 document.getElementById("final").value=document.getElementById("script").value;
};
</script>

</body>
</html>
""")

@app.route("/upload", methods=["POST"])
def upload():
    try:
        return jsonify(url=upload_to_imgbb(request.files["file"]))
    except Exception as e:
        return jsonify(error=str(e)), 500

@app.route("/generate", methods=["POST"])
def generate():
    script = request.form.get("script","")
    result = generate_ppt(script)
    # Handle multiple files
    if isinstance(result, list):
        # Just return first for simplicity
        name, content = result[0]
        return send_file(io.BytesIO(content), download_name=name, as_attachment=True)
    else:
        return send_file(io.BytesIO(result), download_name="output.pptx", as_attachment=True)

app = app
